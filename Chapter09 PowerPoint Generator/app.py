import collections.abc
import config

assert collections
import tkinter as tk
from pptx import Presentation
from pptx.util import Inches, Pt
from openai import OpenAI
from io import BytesIO
import requests

# API Token
client = OpenAI(
  api_key=config.API_KEY,
)

def slide_generator(text, prs):
    prompt = f"Summarize the following text to a DALL-E image generation " \
             f"prompt: \n {text}"

    model_engine = "gpt-4"
    dlp = client.chat.completions.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{prompt}"}
        ],
        max_tokens=250,
        n=1,
        stop=None,
        temperature=0.8
    )

    dalle_prompt = dlp.choices[0].message.content

    response = client.images.generate(
        model="dall-e-3",
        prompt=dalle_prompt + " Style: digital art",
        n=1,
        size="1024x1024"
    )
    image_url = response.data[0].url

    prompt = f"Create a bullet point text for a Powerpoint" \
             f"slide from the following text: \n {text}"
    ppt = client.chat.completions.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{prompt}"}
        ],
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    ppt_text = ppt.choices[0].message.content

    prompt = f"Create a title for a Powerpoint" \
             f"slide from the following text: \n {text}"
    ppt = client.chat.completions.create(
        model=model_engine,
        messages=[
            {"role": "user", "content": "I will ask you a question"},
            {"role": "assistant", "content": "Ok"},
            {"role": "user", "content": f"{prompt}"}
        ],
        max_tokens=1024,
        n=1,
        stop=None,
        temperature=0.8
    )
    ppt_header = ppt.choices[0].message.content

    # Add a new slide to the presentation
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    response = requests.get(image_url)
    img_bytes = BytesIO(response.content)
    slide.shapes.add_picture(img_bytes, Inches(1), Inches(1))

    # Add text box
    txBox = slide.shapes.add_textbox(Inches(3), Inches(1),
                                     Inches(4), Inches(1.5))
    tf = txBox.text_frame
    tf.text = ppt_text

    title_shape = slide.shapes.title
    title_shape.text = ppt_header


def get_slides():
    text = text_field.get("1.0", "end-1c")
    paragraphs = text.split("\n\n")
    prs = Presentation()
    width = Pt(1920)
    height = Pt(1080)
    prs.slide_width = width
    prs.slide_height = height
    for paragraph in paragraphs:
        slide_generator(paragraph, prs)

    prs.save("my_presentation.pptx")


app = tk.Tk()
app.title("Crate PPT Slides")
app.geometry("800x600")

# Create text field
text_field = tk.Text(app)
text_field.pack(fill="both", expand=True)
text_field.configure(wrap="word", font=("Arial", 12))
text_field.focus_set()

# Create the button to create slides
create_button = tk.Button(app, text="Create Slides", command=get_slides)
create_button.pack()

app.mainloop()
